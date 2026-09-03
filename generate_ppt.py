from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(15, 23, 42)
TEXT_WHITE = RGBColor(248, 250, 252)
ACCENT_BLUE = RGBColor(56, 189, 248)
ACCENT_TEAL = RGBColor(45, 212, 191)

def apply_dark_theme(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()

# SLIDE 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide1)
txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.33), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "POTHOLE WALA"
p.font.size = Pt(80)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "AI-Powered Mobile Urban Intelligence Network\nSIH26124"
p2.font.size = Pt(32)
p2.font.color.rgb = TEXT_WHITE
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "\n\nTurning existing public buses into continuously moving urban sensors."
p3.font.size = Pt(24)
p3.font.italic = True
p3.font.color.rgb = ACCENT_TEAL
p3.alignment = PP_ALIGN.CENTER

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

# SLIDE 2
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide2)
add_title(slide2, "POTHOLE WALA - From Road Detection to Verified Repair")
txBox = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "THE PROBLEM"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
for text in ["Road defects are difficult to monitor continuously across large cities.", "Existing complaint systems stop at detection.", "Duplicate reports obscure the physical severity of a defect.", "Municipalities lack actionable, location-based intelligence and repair verification."]:
    p = tf.add_paragraph()
    p.text = "* " + text
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
p = tf.add_paragraph()
p.text = "\nOUR SOLUTION"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
for text in ["1. SENSE: Existing public buses capture road conditions during regular routes.", "2. DETECT: Edge AI identifies defects and attaches GPS/time evidence.", "3. INTELLIGENTLY ACT: Multiple observations are spatially fused, prioritized, and ticketed.", "4. VERIFY: A later bus pass re-inspects the location and marks it VERIFIED or REOPENED."]:
    p = tf.add_paragraph()
    p.text = "  " + text
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
p = tf.add_paragraph()
p.text = "\n> Most systems detect. POTHOLE WALA detects -> acts -> verifies."
p.font.size = Pt(24)
p.font.bold = True
p.font.italic = True
p.font.color.rgb = ACCENT_BLUE

# SLIDE 3
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide3)
add_title(slide3, "Technical Architecture & Methodology")
txBox = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.add_paragraph()
p.text = "CLOSED-LOOP VERIFICATION FLOW"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
flow = "BUS CAMERA -> EDGE AI -> TRACKING -> DETECTION EVENT\n|\nV\nFASTAPI -> POSTGRESQL + POSTGIS (Spatial Fusion) -> URBAN ISSUE -> PRIORITY ENGINE\n|\nV\nMUNICIPAL TICKET -> REPAIR REPORTED -> BUS REVISIT -> AI VERIFICATION -> VERIFIED / REOPENED"
p = tf.add_paragraph()
p.text = flow
p.font.size = Pt(18)
p.font.color.rgb = TEXT_WHITE
p = tf.add_paragraph()
p.text = "\nTECHNOLOGIES"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
for t in ["Frontend: React, TypeScript, Vite", "Backend: FastAPI, Server-Sent Events (SSE)", "Database: PostgreSQL, PostGIS (ST_DWithin)", "AI: YOLOv8-compatible pipeline, Centroid Tracking", "Deployment: Docker Compose"]:
    p = tf.add_paragraph()
    p.text = "* " + t
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
p = tf.add_paragraph()
p.text = "\n*Note: Prototype currently uses deterministic fallback inference for reliable demonstration."
p.font.size = Pt(14)
p.font.italic = True
p.font.color.rgb = RGBColor(156, 163, 175)

# SLIDE 4
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide4)
add_title(slide4, "Feasibility, Risks & Mitigation")
txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
rows = [
    ("Real-time processing", "Edge-first inference + event-based dispatch", "TensorRT edge models"),
    ("Duplicate detections", "Temporal tracking + PostGIS spatial fusion", "Advanced confidence fusion"),
    ("Poor connectivity", "Event-based HTTP dispatch", "Offline edge SQLite buffer"),
    ("Trusting repairs", "Independent bus revisit + AI verification", "Automated payouts")
]
p = tf.add_paragraph()
p.text = "CHALLENGE                 | MITIGATION                                      | FUTURE SCALE"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT_TEAL
for r in rows:
    p = tf.add_paragraph()
    p.text = f"{r[0].ljust(25)} | {r[1].ljust(47)} | {r[2]}"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_WHITE
p = tf.add_paragraph()
p.text = "\nVIABILITY PILOT: 1 City -> Selected Routes -> Municipal Dashboard -> Expand Fleet."
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

# SLIDE 5
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide5)
add_title(slide5, "Impact - From Pothole Detection to Accountable Road Maintenance")
txBox = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
impacts = [
    ("1. CITIZENS", "Faster identification of hazards, better road safety."),
    ("2. MUNICIPALITIES", "Location-based intelligence, prioritized maintenance, reduced duplicate complaints."),
    ("3. PUBLIC FUNDS", "Traceable repair actions, independent verification, failed repairs automatically reopened."),
    ("4. CITY INFRASTRUCTURE", "Existing buses become mobile sensing assets. Expandable to signs and markings.")
]
for title, desc in impacts:
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "  " + desc
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p = tf.add_paragraph()
p = tf.add_paragraph()
p.text = "> Every bus journey can become another observation of the city's road network."
p.font.size = Pt(24)
p.font.bold = True
p.font.italic = True
p.font.color.rgb = ACCENT_BLUE

# SLIDE 6
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
apply_dark_theme(slide6)
add_title(slide6, "Research, Technologies & References")
txBox = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
tf = txBox.text_frame
tf.word_wrap = True
refs = [
    "AI & Computer Vision", "  * Ultralytics YOLOv8 Architecture Documentation", "  * Object Tracking & Duplicate Suppression Methodologies (Centroid/Euclidean)",
    "Geospatial & Database", "  * PostgreSQL 15 & PostGIS 3.3 Documentation", "  * Spatial Indexing and ST_DWithin geofencing methodologies",
    "Backend & Real-Time", "  * FastAPI / Starlette Documentation", "  * Server-Sent Events (SSE) W3C Specifications",
    "Urban Mobility", "  * Intelligent Transportation Systems (ITS) Smart City frameworks"
]
for line in refs:
    p = tf.add_paragraph()
    p.text = line
    if line.startswith("  *"):
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_WHITE
    else:
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL

prs.save('POTHOLE_WALA_SIH_FINAL.pptx')
print('Successfully saved POTHOLE_WALA_SIH_FINAL.pptx')
